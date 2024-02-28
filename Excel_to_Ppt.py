from pptx import Presentation
import pandas as pd
from math import ceil
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.dml.color import RGBColor  # Corrected import
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Inches


def process_files(adscore_excel, example_excel_url, ppt_file,selected_option):
    # ppt_file_path = 'AdScore Reporting Template.pptx'
    prs = Presentation(ppt_file)

    example_excel = pd.ExcelFile(example_excel_url)
    example_csv= "Hotel_Example.csv"
    df = pd.read_excel(example_excel, 'CAMP15_0902_v5', header=None)
    df.to_csv(example_csv, index=False, header=False)
    df = pd.read_csv(example_csv, header=None)

    # AdScore Norms W46 (W87)

    AdScore_excel = pd.ExcelFile(adscore_excel)
    AdScore_csv= "AdScore.csv"
    AdScore_df = pd.read_excel(AdScore_excel, 'Data', header=None)
    AdScore_df.to_csv(AdScore_csv, index=False, header=False)
    AdScore_df = pd.read_csv(AdScore_csv, header=None)

    # Example Brand and compaign
    example_excel = pd.ExcelFile(example_excel_url)
    example_csv_brand_compaign= "Example_for_brand_compaign.csv"
    example_df_for_brand_compaign = pd.read_excel(example_excel, 'TOC', header=None)
    example_df_for_brand_compaign.to_csv(example_csv_brand_compaign, index=False, header=False)
    example_df_for_brand_compaign = pd.read_csv(example_csv_brand_compaign, header=None)


    # Airline Data
    # airline_excel = pd.ExcelFile('Airlines example.xlsx')
    # airline_csv = "Example_for_brand_compaign.csv"
    # Airline_df = pd.read_excel(airline_excel, 'CAMP15_0902_v5', header=None)
    # Airline_df.to_csv(airline_csv, index=False, header=False)
    # Airline_df = pd.read_csv(airline_csv, header=None)

    # In[5]:

    try:
        slide_index = 1 
        # map_value1= calc* = (XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C131 / XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D131) -1
        # Convert the values to numeric (float for safety in division)
        map_value1 = f"{ceil((pd.to_numeric(df.iloc[130, 2], errors='coerce') / pd.to_numeric(df.iloc[130, 3], errors='coerce') - 1) * 100)}%"
        # map_value2 Yellow Highlight: calc* = (XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C262 / XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D262) -1
        map_value2 = f"{ceil((pd.to_numeric(df.iloc[261, 2], errors='coerce') / pd.to_numeric(df.iloc[261, 3], errors='coerce') - 1) * 100)}%"
        # map_value3 Yellow Highlight: calc* = (XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C281 / XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D281) -1
        map_value3 = f"{ceil((pd.to_numeric(df.iloc[280, 2], errors='coerce') / pd.to_numeric(df.iloc[280, 3], errors='coerce') - 1) * 100)}%"
        # map_value4 Yellow Highlight: calc* = (XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C148 / XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D148) -1
        map_value4 = f"{ceil((pd.to_numeric(df.iloc[147, 2], errors='coerce') / pd.to_numeric(df.iloc[147, 3], errors='coerce') - 1) * 100)}%"

        replacement_values = [map_value1,map_value3, map_value2,map_value4]

        replacement_index = 0

        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]

            for index, shape in enumerate(slide.shapes):
                text=f'The campaign was particularly effective in driving knowledge (familiarity) & recommendation for {example_df_for_brand_compaign.iloc[0, 1]}'
                if index == 10:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index ==0 :
                            if paragraph.runs:
                                paragraph.runs[0].text = text
                                for run in paragraph.runs[1:]:
                                    run.text = ""  


            def process_and_replace_shape_with_formatting(shape, values, parent_index=""):
                nonlocal replacement_index  # This tells Python to use the variable from the enclosing scope
                if shape.shape_type == 6:  
                    for i, sub_shape in enumerate(shape.shapes):
                        composite_index = f"{parent_index}-{i}" if parent_index else str(i)
                        process_and_replace_shape_with_formatting(sub_shape, values, composite_index)
                else:
                    if shape.has_text_frame and "%" in shape.text_frame.text and replacement_index < len(values):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if "%" in run.text:
                                    # Preserve formatting but replace tex
                                    if replacement_index < len(values):
                                        new_text = str(values[replacement_index])
                                        run.text = new_text
                                        replacement_index += 1
                                        break  # Assuming only one replacement per shape

            for index, shape in enumerate(slide.shapes):
                process_and_replace_shape_with_formatting(shape, replacement_values, str(index))
            print("----- Slide No 1 Completed")
            prs.save('AdScore Reporting Template.pptx')
    except Exception as e:
            print(f"An error occurred at slide 1 :  {e}")


    # # Slide 4

    # In[6]:


    try:
        slide_index = 3
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            # making market as per c12-c39 data 
            markets_data = "" 
            for i in range(11, 39):  
                if pd.to_numeric(df.iloc[i, 2], errors='coerce') > 0:  
                    markets_data += str(df.iloc[i, 0]) + ","  
            new_data_for_source = df.iloc[0, 0]
            new_data_for_source =f'Source: {new_data_for_source}'
            test_value =  pd.to_numeric(df.iloc[10, 2], errors='coerce')  #C11
            control_value = pd.to_numeric(df.iloc[10, 3], errors='coerce') # D11
            new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table  
                    for row_index, row in enumerate(table.rows):
                        for row_index, row in enumerate(table.rows):
                        # Check the condition for the first cell of each row to decide the updates
                            first_cell_text = row.cells[0].text.strip()
                        # Determine the updates based on row_index and the content of the first cell
                            updates = {}
                            if row_index == 2 and first_cell_text == "Male":
                                updates = {1: str(ceil(pd.to_numeric(df.iloc[70, 3], errors='coerce')*100))+ "%", 
                                            2:str(ceil(pd.to_numeric(df.iloc[70, 2], errors='coerce')*100))+ "%"  }
                            elif row_index == 3 and first_cell_text == "Female":
                                updates = {1: str(ceil(pd.to_numeric(df.iloc[71, 3], errors='coerce')*100))+ "%", 
                                            2:str(ceil(pd.to_numeric(df.iloc[71, 2], errors='coerce')*100))+ "%"  }
                            elif row_index ==5 and first_cell_text  == "18-34 yrs":
                                updates = {1: str(ceil(pd.to_numeric(df.iloc[89, 3], errors='coerce')*100))+ "%", 
                                            2:str(ceil(pd.to_numeric(df.iloc[89, 2], errors='coerce')*100))+ "%"  }
                            elif row_index ==6 and (first_cell_text  == "35-54 yrs" or first_cell_text == "35-44 yrs"):
                                updates = {1: str(ceil(pd.to_numeric(df.iloc[90, 3], errors='coerce')*100))+ "%", 
                                            2:str(ceil(pd.to_numeric(df.iloc[90, 2], errors='coerce')*100))+ "%"  }
                            elif row_index ==7 and (first_cell_text  == "55+ yrs" or first_cell_text == "45+ yrs"):
                                updates = {1: str(ceil(pd.to_numeric(df.iloc[91, 3], errors='coerce')*100))+ "%", 
                                            2:str(ceil(pd.to_numeric(df.iloc[91, 2], errors='coerce')*100))+ "%"  }
                            for cell_index, new_text in updates.items():
                                cell = row.cells[cell_index]
                                if cell.text_frame:
                                    if cell.text_frame.paragraphs:
                                        paragraph = cell.text_frame.paragraphs[0]
                                        if paragraph.runs:
                                            paragraph.runs[0].text = new_text
                                        else:
                                            paragraph.add_run(new_text)
                                    else:
                                        # Create a new text_frame if not present (unlikely for table cells but just in case)
                                        cell.text = new_text  # This will replace all content and potentially formatting
                                else:
                                    # Directly setting text if no text frame is found, fallback option
                                    cell.text = new_text

            for index, shape in enumerate(slide.shapes):
                shape_type = shape.shape_type
                text = shape.text if shape.has_text_frame else "No text"
                if index == 10:
                    if shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text = "Markets:"
                        paragraph.runs[1].text = markets_data
                if  index == 3:
                    if shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
            #                 replacing Source: Adscore Wave 84 12th – 25th June 2023
                            if paragraph_index== 3:
                                paragraph.text = new_data_for_source
            #                   Sample Test n=104 , Control n=52 *Caution, low base size                
                            if paragraph_index== 4:
                                paragraph.text = new_data_for_source1                    


            print("----- Slide No 4 Completed")

            prs.save('AdScore Reporting Template.pptx')
    except Exception as e:
        print(f"An error occurred at slide 4:  {e}")    


    # # Slide 6

    # In[7]:


    try:
        slide_index = 5
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            chart = None

            new_data_for_source = df.iloc[0, 0]
            new_data_for_source =f'Source: {new_data_for_source}'
            test_value =  pd.to_numeric(df.iloc[10, 2], errors='coerce')  #C11
            control_value = pd.to_numeric(df.iloc[10, 3], errors='coerce') # D11
            new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'
            new_value_for_pie_chart= ceil(pd.to_numeric(df.iloc[350, 2], errors='coerce')) #C351
            new_value_for_Tv1 = str(ceil(pd.to_numeric(df.iloc[335, 2], errors='coerce')*100))+ "%"  #C336
            new_value_for_Tv2 = str(ceil(pd.to_numeric(df.iloc[320, 2], errors='coerce')*100))+ "%" #C321
            #Yellow Benchmark Highlight: calc* = “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: DV8)
            Bencmark1 =str(ceil(pd.to_numeric(AdScore_df.iloc[7, 125], errors='coerce')*100))+ "%"
            Bencmark1=f'BENCHMARK: {Bencmark1}'
            # “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: DZ8)
            Bencmark2 =str(ceil(pd.to_numeric(AdScore_df.iloc[7, 129], errors='coerce')*100))+ "%"
            Bencmark2=f'BENCHMARK: {Bencmark2}'
            # value="254%"
            # compaign_Recall_were=f'Those most likely to recall the campaign were: {value}'

            # pie chart update
            another_value=100-new_value_for_pie_chart
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type == XL_CHART_TYPE.PIE:
                    chart = shape.chart
                    break

            if chart:
                # Create a new CategoryChartData object
                chart_data = CategoryChartData()

                # Assuming we want to set the first segment to 60% and distribute the rest
                new_values = [another_value, new_value_for_pie_chart]  # Example values, adjust based on your actual data needs

                # Update the chart data
                chart_data.categories = ['Category 1', 'Category 2']  # Update categories as needed
                chart_data.add_series('Series 1', new_values)

                # Replace existing chart data
                chart.replace_data(chart_data)

            for index, shape in enumerate(slide.shapes):
                text=f'Those most likely to recall the campaign were:{example_df_for_brand_compaign.iloc[1, 1]}'
                if index == 2:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index ==4 :
                            if paragraph.runs:
                                paragraph.runs[0].text = text
                                for run in paragraph.runs[1:]:
                                    run.text = "" 

                if index == 6 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):         
                        if paragraph_index ==2:
                            paragraph.text = new_data_for_source           
                        if paragraph_index ==3:
                            paragraph.text =new_data_for_source1           
                if index == 9 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text=str(new_value_for_pie_chart)+"%"
                if index == 11 and shape.has_text_frame: #C336
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text = new_value_for_Tv1

                if index == 12 and shape.has_text_frame:  # C321
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text=new_value_for_Tv2

                if index == 15 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index ==0 :
                            if paragraph.runs:
                                paragraph.runs[0].text = Bencmark2
                                for run in paragraph.runs[1:]:
                                    run.text = ""            
                if index == 16 and shape.has_text_frame:   #Yellow Benchmark Highlight: calc* = “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: DV8) 
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index ==0 :
                            if paragraph.runs:
                                paragraph.runs[0].text = Bencmark1
                                for run in paragraph.runs[1:]:
                                    run.text = ""  
            print("----- Slide No 6 Completed")
            prs.save('AdScore Reporting Template.pptx')
    except Exception as e:
        print(f"An error occurred at slide 6:  {e}") 


    # # Slide 7

    # In[8]:


    try:
        slide_index = 6
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            new_data_for_source = df.iloc[0, 0]
            new_data_for_source =f'Source: {new_data_for_source}'
            sample_size_cell = pd.to_numeric(df.iloc[380, 2], errors='coerce') # C381
            source_adscoorce_line2 = f'Sample Test All those who recalled advertising n={sample_size_cell} *Caution, low base size'
            for index, shape in enumerate(slide.shapes):
                if index == 12 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= str(ceil(pd.to_numeric(df.iloc[615, 2], errors='coerce')*100))+ "%"  #C616
                if index == 15 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= str(ceil(pd.to_numeric(df.iloc[619, 2], errors='coerce')*100))+ "%" #C620
                if index == 18 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= str(ceil(pd.to_numeric(df.iloc[617, 2], errors='coerce')*100))+ "%" #C618
                if index == 21 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text=  str(ceil(pd.to_numeric(df.iloc[620, 2], errors='coerce')*100))+ "%" #C621
                if index == 24 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= str(ceil(pd.to_numeric(df.iloc[618, 2], errors='coerce')*100))+ "%"  #C619
                if index == 27 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= str(ceil(pd.to_numeric(df.iloc[621, 2], errors='coerce')*100))+ "%"  #C622
            #         _________________________________________________________________________

            # Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  C616 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: EF8) -1
                if index == 14 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= f"{ceil((pd.to_numeric(df.iloc[615, 2], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 135], errors='coerce') - 1) * 100)}% "
            # Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  C620 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: EP8) -1
                if index == 17 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text=f"{ceil((pd.to_numeric(df.iloc[619, 2], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 145], errors='coerce') - 1) * 100)}% "
            # Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  C618 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: EL8) -1
                if index == 20 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= f"{ceil((pd.to_numeric(df.iloc[617, 2], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 141], errors='coerce') - 1) * 100)}% "
            # Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  C621 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: ER8) -1
                if index == 23 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= f"{ceil((pd.to_numeric(df.iloc[620, 2], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 147], errors='coerce') - 1) * 100)}% vs Benchmarks "
                    for run in paragraph.runs[1:]:
                        run.text = "" 
            # Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  C619 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: EN8) -1
                if index == 26 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= f"{ceil((pd.to_numeric(df.iloc[618, 2], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 143], errors='coerce') - 1) * 100)}% "
            # Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  C622 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: ET8) -1
                if index == 29 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= f"{ceil((pd.to_numeric(df.iloc[621, 2], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 149], errors='coerce') - 1) * 100)}% vs Benchmarks "
                    for run in paragraph.runs[1:]:
                        run.text = "" 
            #         -----------------------------------------------------------------------------
                if index == 3 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index == 1:
                            paragraph.runs[0].text = new_data_for_source
                            for run in paragraph.runs[1:]:
                                run.text = "" 
                        if paragraph_index == 2:
                            paragraph.runs[0].text = source_adscoorce_line2
                            for run in paragraph.runs[1:]:
                                run.text = "" 
            print("----- Slide No 7 Completed")
            prs.save('AdScore Reporting Template.pptx')
    except Exception as e:
        print(f"An error occurred at slide 7:  {e}") 


    # # Slide 8

    # In[9]:


    try:
        slide_index = 7  
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
        new_data_for_source = df.iloc[0, 0]
        new_data_for_source =f'Source: {new_data_for_source}'
        sample_size_cell = pd.to_numeric(df.iloc[10, 2], errors='coerce') # C11
        source_adscoorce_line2 = f'Sample Test All those who recalled advertising n={sample_size_cell} *Caution, low base size'
        for index, shape in enumerate(slide.shapes):
            if  index == 4:
                if shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index== 0:
                            paragraph.runs[0].text =  str(ceil(pd.to_numeric(df.iloc[616, 2], errors='coerce')*100))+ "%"  #C617
            if  index == 8:
                if shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
        # Yellow Benchmark Highlight: calc* = “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: EJ8) 
                        if paragraph_index== 0:
                            paragraph.runs[0].text = str(ceil(pd.to_numeric(AdScore_df.iloc[7, 139], errors='coerce')*100))+ "%"  # Adscore EJ8 
        #            Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  C617 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: EJ8) -1
                        if paragraph_index== 2:
                            paragraph.runs[0].text = f"{ceil((pd.to_numeric(df.iloc[616, 2], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 139], errors='coerce') - 1) * 100)}% vs"
                            for run in paragraph.runs[1:]:
                                run.text = "" 
            if index == 3 and shape.has_text_frame:
                paragraph = shape.text_frame.paragraphs[0]
                for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                    if paragraph_index == 1:
                        paragraph.runs[0].text = new_data_for_source
                        for run in paragraph.runs[1:]:
                            run.text = "" 
                    if paragraph_index == 2:
                        paragraph.runs[0].text = source_adscoorce_line2
                        for run in paragraph.runs[1:]:
                            run.text = "" 
        print("----- Slide No 8 Completed")
        prs.save('AdScore Reporting Template.pptx')
    except Exception as e:
        print(f"An error occurred at slide 8:  {e}") 


    # # Slide 10

    # In[10]:
    try:
        slide_index = 9  
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            new_data_for_source = df.iloc[0, 0]
            new_data_for_source =f'Source: {new_data_for_source}'
            test_value =  pd.to_numeric(df.iloc[10, 2], errors='coerce')  #C11
            control_value = pd.to_numeric(df.iloc[10, 3], errors='coerce') # D11
            new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'

            for index, shape in enumerate(slide.shapes):
                if  index == 7:
                    if shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
            # Yellow Highlight: calc* = (XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C131 / XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D131) -1
                            if paragraph_index== 0:
                                paragraph.runs[0].text = f"{ceil((pd.to_numeric(df.iloc[130, 2], errors='coerce') / pd.to_numeric(df.iloc[130, 3], errors='coerce') - 1) * 100)}% "
                                for run in paragraph.runs[1:]:
                                    run.text = "" 
                if index == 4 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index == 1:
                            paragraph.runs[0].text = new_data_for_source
                            for run in paragraph.runs[1:]:
                                run.text = "" 
                        if paragraph_index == 2:
                            paragraph.runs[0].text = new_data_for_source1
                            for run in paragraph.runs[1:]:
                                run.text = "" 
                if index == 10:
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
            #Yellow Benchmark Highlight: calc* = “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: Z8) 
                        if paragraph_index == 0:
                            paragraph.runs[0].text = str(ceil(pd.to_numeric(AdScore_df.iloc[7, 25], errors='coerce')*100))+ "%"  #Z8
                            for run in paragraph.runs[1:]:
                                run.text = "" 
            #Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  C131 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: Z8) -1
                        if paragraph_index == 2:
                            paragraph.runs[0].text =  f"{ceil((pd.to_numeric(df.iloc[130, 2], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 25], errors='coerce') - 1) * 100)}% vs"
                            for run in paragraph.runs[1:]:
                                run.text = "" 
            bar_value1=ceil(pd.to_numeric(df.iloc[130, 3], errors='coerce')*100) / 100#D131
            bar_value2=ceil(pd.to_numeric(df.iloc[130, 2], errors='coerce')*100) / 100 # C131 

            chart_data = CategoryChartData()
            chart_data.categories = ['Total Brand Awareness']  
            chart_data.add_series('Series 2', (bar_value1,)) 
            chart_data.add_series('Series 1', (bar_value2,)) 
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
                    chart = shape.chart
                    chart.replace_data(chart_data)

                    # Apply formatting to each series in the chart
                    for series in chart.series:
                        series.data_labels.number_format = '0%'  # Format as percentage
                        series.data_labels.show_value = True     # Ensure values are shown

                    break
            print("----- Slide No 10 Completed")
            prs.save('AdScore Reporting Template.pptx')

    except Exception as e:
        print(f"An error occurred at slide 10:  {e}")                


    # # Slide 11

    # In[11]:


    try:
        slide_index = 10 
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            new_data_for_source = df.iloc[0, 0]
            new_data_for_source =f'Source: {new_data_for_source}'
            test_value =  pd.to_numeric(df.iloc[10, 2], errors='coerce')  #C11
            control_value = pd.to_numeric(df.iloc[10, 3], errors='coerce') # D11
            new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'
            control1= ceil(pd.to_numeric(df.iloc[124, 3], errors='coerce')*100)  #D125
            control1 =control1 /100
            control2= ceil(pd.to_numeric(df.iloc[125, 3], errors='coerce')*100)  #D126
            control2 =control2 /100
            test1= ceil(pd.to_numeric(df.iloc[124, 2], errors='coerce')*100)  #C125
            test1 =test1 /100
            test2= ceil(pd.to_numeric(df.iloc[125, 2], errors='coerce')*100)   #C126
            test2 =test2 /100

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type == XL_CHART_TYPE.BAR_STACKED:
                    chart = shape.chart
                    # Assuming the chart has exactly two series and you're updating both
                    new_values_series_1 = [control1, test1]  # New values for the first series
                    new_values_series_2 = [control2, test2]  # New values for the second series
                    chart_data = CategoryChartData()
                    chart_data.categories = ['CONTROL', 'TEST'] 
                    chart_data.add_series('Series 1', new_values_series_1)
                    chart_data.add_series('Series 2', new_values_series_2)
                    for series in chart.series:
                        series.data_labels.number_format = '0%'
                    # Replace the old data with the new data
                    chart.replace_data(chart_data)
                    break

            for index, shape in enumerate(slide.shapes):
            #     D131
                if index == 8 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= str(ceil(pd.to_numeric(df.iloc[130, 3], errors='coerce')*100))+ "%"  
            #         C131
                if index == 9 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= str(ceil(pd.to_numeric(df.iloc[130, 2], errors='coerce')*100))+ "%"  
            #     Yellow Highlight: calc* = (XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C131 / XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D131) -1
                if index == 10:
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index == 0:
                            paragraph.runs[0].text = f"{ceil((pd.to_numeric(df.iloc[130, 2], errors='coerce') / pd.to_numeric(df.iloc[130, 3], errors='coerce') - 1) * 100)}%"
                if index == 14:
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
            # Yellow Benchmark Highlight: calc* = “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: Z8) 
                        if paragraph_index == 0:
                            paragraph.runs[0].text = str(ceil(pd.to_numeric(AdScore_df.iloc[7, 25], errors='coerce')*100))+ "%"  #Z8
                            for run in paragraph.runs[1:]:
                                run.text = "" 
            # Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  C131 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: Z8) -1
                        if paragraph_index == 2:
                            paragraph.runs[0].text =  f"{ceil((pd.to_numeric(df.iloc[130, 2], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 25], errors='coerce') - 1) * 100)}% vs"
                            for run in paragraph.runs[1:]:
                                run.text = "" 
                if index == 5 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index == 1:
            #                 Source cell: A1
                            paragraph.runs[0].text = new_data_for_source
                            for run in paragraph.runs[1:]:
                                run.text = "" 
                        if paragraph_index == 2:
            #                 Test sample size cell: C11 Control sample size cell: D11
                            paragraph.runs[0].text = new_data_for_source1
                            for run in paragraph.runs[1:]:
                                run.text = ""                 
            print("----- Slide No 11 Completed")
            prs.save('AdScore Reporting Template.pptx')
    except Exception as e:
        print(f"An error occurred at slide 11:  {e}")                 


    # # Slide 12

    # In[12]:


    try:
        slide_index = 11
        if slide_index < len(prs.slides):
            Airline_df = df
            slide = prs.slides[slide_index]
            new_data_for_source = df.iloc[0, 0]
            new_data_for_source =f'Source: {new_data_for_source}'
            test_value =  pd.to_numeric(df.iloc[10, 2], errors='coerce')  #C11
            control_value = pd.to_numeric(df.iloc[10, 3], errors='coerce') # D11
            new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'
            # Grey bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D262
            # Red bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C262
            bar_value1 = ceil(pd.to_numeric(Airline_df.iloc[261, 3], errors='coerce')*100)
            bar_value1 = bar_value1/100
            bar_value2 =ceil(pd.to_numeric(Airline_df.iloc[261, 2], errors='coerce')*100)
            bar_value2 =bar_value2 / 100
            # Grey bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D281
            #Red bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C281
            bar_value3 =ceil(pd.to_numeric(Airline_df.iloc[280, 3], errors='coerce')*100)
            bar_value3 =bar_value3/100
            bar_value4 =ceil(pd.to_numeric(Airline_df.iloc[280, 2], errors='coerce')*100)
            bar_value4= bar_value4 /100

            # Grey bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D148
            # Red bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C148
            bar_value5 =ceil(pd.to_numeric(Airline_df.iloc[147, 3], errors='coerce')*100)
            bar_value5 =bar_value5 /100
            bar_value6 =ceil(pd.to_numeric(Airline_df.iloc[147, 2], errors='coerce')*100)
            bar_value6 =bar_value6/100

            for index, shape in enumerate(slide.shapes):
            # Yellow Highlight: calc* = (XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C262 / XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D262) -1

                if index == 1 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= f"{ceil((pd.to_numeric(df.iloc[261, 2], errors='coerce') / pd.to_numeric(df.iloc[261, 3], errors='coerce') - 1) * 100)}"
            # Yellow Highlight: calc* = (XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C281 / XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D281) -1
                if index == 3 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= f"{ceil((pd.to_numeric(df.iloc[280, 2], errors='coerce') / pd.to_numeric(df.iloc[280, 3], errors='coerce') - 1) * 100)}"
            #     Yellow Highlight: calc* = (XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C148 / XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D148) -1
                if index == 13 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    paragraph.runs[0].text= f"{ceil((pd.to_numeric(df.iloc[147, 2], errors='coerce') / pd.to_numeric(df.iloc[147, 3], errors='coerce') - 1) * 100)}"    
                if index == 9 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index == 3:
            #                 Source cell: A1
                            paragraph.runs[0].text = new_data_for_source
                            for run in paragraph.runs[1:]:
                                run.text = "" 
                        if paragraph_index == 4:
            #                 Test sample size cell: C11 Control sample size cell: D11
                            paragraph.runs[0].text = new_data_for_source1
                            for run in paragraph.runs[1:]:
                                run.text = ""                 
                if index == 17 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index == 0:
            # Yellow Bechmark Highlight: calc* = “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: DJ8) 
                            paragraph.runs[0].text = str(ceil(pd.to_numeric(AdScore_df.iloc[7, 113], errors='coerce')*100))+ "% Benchmark" 
                            for run in paragraph.runs[1:]:
                                run.text = "" 
                        if paragraph_index == 1:
            # Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D262 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: DJ8) -1
                            paragraph.runs[0].text =  f"{ceil((pd.to_numeric(df.iloc[261, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 113], errors='coerce') - 1) * 100)}% Benchmark" 
                            for run in paragraph.runs[1:]:
                                run.text = ""     

                if index == 18 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index == 0:
            # Yellow Bechmark Highlight: calc* = “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: DN8) 
                            paragraph.runs[0].text = str(ceil(pd.to_numeric(AdScore_df.iloc[7, 117], errors='coerce')*100))+ "% Benchmark" 
                            for run in paragraph.runs[1:]:
                                run.text = "" 
                        if paragraph_index == 1:
            # Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D281 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: DN8) -1
                            paragraph.runs[0].text =  f"{ceil((pd.to_numeric(df.iloc[280, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 117], errors='coerce') - 1) * 100)}% Benchmark" 
                            for run in paragraph.runs[1:]:
                                run.text = ""     

                if index == 19 and shape.has_text_frame:
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index == 0:
            # Yellow Bechmark Highlight: calc* = “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AD8) 
                            paragraph.runs[0].text = str(ceil(pd.to_numeric(AdScore_df.iloc[7, 29], errors='coerce')*100))+ "% Benchmark" 
                            for run in paragraph.runs[1:]:
                                run.text = "" 
                        if paragraph_index == 1:
            # Yellow Highlight: calc*=  (File “XXX Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D148 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AD8) -1
                            paragraph.runs[0].text =  f"{ceil((pd.to_numeric(df.iloc[147, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 29], errors='coerce') - 1) * 100)}% Benchmark" 
                            for run in paragraph.runs[1:]:
                                run.text = ""     


            chart_data = CategoryChartData()
            chart_data.categories = ['Brand Recomendation', 'Brand consideration','Brand Sentiment']  
            chart_data.add_series('Series 1', (bar_value1, bar_value3, bar_value5))  # Adjust these values as needed
            chart_data.add_series('Series 2', (bar_value2, bar_value4, bar_value6))  # Adjust these values as needed
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type ==XL_CHART_TYPE.COLUMN_CLUSTERED:
                    chart = shape.chart
                    chart.replace_data(chart_data)

                    # Apply formatting to each series in the chart
                    for series in chart.series:
                        series.data_labels.number_format = '0%'  # Format as percentage
                        series.data_labels.show_value = True     # Ensure values are shown

                    break
            print("----- Slide No 12 Completed")
            prs.save('AdScore Reporting Template.pptx')
    except Exception as e:
        print(f"An error occurred at slide 12:  {e}")    


# ---------------------------------------------LAST Slide-----------------------------------------------------------------------------------#  
    # # Slide 21 Lastslide
    # In[29]:
    try:
        # Check if the slide index exists
        slide_index = 20  # Example slide index
        if slide_index < len(prs.slides):
            other_df = df
            slide = prs.slides[slide_index]
            for index, shape in enumerate(slide.shapes):
                if index == 1 :
                    paragraph = shape.text_frame.paragraphs[0]
                    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                        if paragraph_index == 0:
                            # Cell: C351
                            paragraph.runs[0].text = "Good recall and cut through at "+str(ceil(pd.to_numeric(df.iloc[350, 2], errors='coerce')*100))+ "%" 
                            for run in paragraph.runs[1:]:
                                run.text = "" 
            #                     Cell: C616
                        if paragraph_index == 3:
                            paragraph.runs[0].text = str(ceil(pd.to_numeric(df.iloc[615, 2], errors='coerce')*100))+ "% Got my attention & Aimed at people like me" 
                            for run in paragraph.runs[1:]:
                                run.text = "" 
            # Cell: C617
                        if paragraph_index == 6:
                            paragraph.runs[0].text = str(ceil(pd.to_numeric(df.iloc[616, 2], errors='coerce')*100))+ "% Positive fit" 
                            for run in paragraph.runs[1:]:
                                run.text = "" 
            #Yellow Highlight: calc* = (XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C131 / XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D131) -1
                        if paragraph_index == 9:
                            paragraph.runs[0].text = f"{ceil((pd.to_numeric(df.iloc[130, 2], errors='coerce') / pd.to_numeric(other_df.iloc[130, 3], errors='coerce') - 1) * 100)}% uplift on the metric" 
                            for run in paragraph.runs[1:]:
                                run.text = ""     
            #                     Yellow Highlight: calc* = (XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C262 / XXX Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D262) -1
                        if paragraph_index == 12:
                            paragraph.runs[0].text = f"{ceil((pd.to_numeric(df.iloc[261, 2], errors='coerce') / pd.to_numeric(other_df.iloc[261, 3], errors='coerce') - 1) * 100)}% uplift between control & test" 
                            for run in paragraph.runs[1:]:
                                run.text = ""     
            prs.save('AdScore Reporting Template.pptx')
    except Exception as e:
        print(f"An error occurred at slide 21:  {e}") 

    # # Slide 13

    # In[13]:

    if selected_option == "Airline":
        # Airline Data for slide 12 
        # Airline_excel = pd.ExcelFile('Airlines example.xlsx')
        # Airline_csv= "Airline.csv"
        # Airline_df = pd.read_excel(Airline_excel, 'CAMP15_0902_v5', header=None)
        # Airline_df.to_csv(Airline_csv, index=False, header=False)
        # Airline_df = pd.read_csv(Airline_csv, header=None)
        Airline_df = df
        try:
            slide_index = 12 
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                new_data_for_source = Airline_df.iloc[0, 0]
                new_data_for_source =f'Source: {new_data_for_source}'
                test_value =  pd.to_numeric(Airline_df.iloc[10, 2], errors='coerce')  #C11
                control_value = pd.to_numeric(Airline_df.iloc[10, 3], errors='coerce') # D11
                new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'

                # Grey bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243
                # Red bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243
                # Yellow Highlight: calc*=  (File “Airlines Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D243 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AN8)-1
                market_leader1 = ceil(pd.to_numeric(Airline_df.iloc[242, 3], errors='coerce')*100) / 100
                market_leader2 = ceil(pd.to_numeric(Airline_df.iloc[242, 2], errors='coerce')*100) / 100
                bencmark1= f"{ceil((pd.to_numeric(Airline_df.iloc[242, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 39], errors='coerce') - 1) * 100)}%"

                # Grey bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224
                # Red bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224
                # Yellow Highlight: calc*=  (File “Airlines Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D224 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AL8) -1
                connect_world1 = ceil(pd.to_numeric(Airline_df.iloc[223, 3], errors='coerce')*100) / 100
                connect_world2 = ceil(pd.to_numeric(Airline_df.iloc[223, 2], errors='coerce')*100) / 100
                bencmark2 = f"{ceil((pd.to_numeric(Airline_df.iloc[223, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,37 ], errors='coerce') - 1) * 100)}%"

                # Grey bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205
                # Red bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205
                # Yellow Highlight: calc*=  (File “Airlines Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D205 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AJ8) -1
                brand_love1 = ceil(pd.to_numeric(Airline_df.iloc[204, 3], errors='coerce')*100) / 100
                brand_love2 = ceil(pd.to_numeric(Airline_df.iloc[204, 2], errors='coerce')*100) / 100
                bencmark3 = f"{ceil((pd.to_numeric(Airline_df.iloc[204, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,35 ], errors='coerce') - 1) * 100)}%"

                # Grey bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186
                # Red bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186
                # Yellow Highlight: calc*=  (File “Airlines Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D186 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AH8) -1
                flight_exp1 = ceil(pd.to_numeric(Airline_df.iloc[185, 3], errors='coerce')*100) / 100
                flight_exp2 = ceil(pd.to_numeric(Airline_df.iloc[185, 2], errors='coerce')*100) / 100
                bencmark4 = f"{ceil((pd.to_numeric(Airline_df.iloc[185, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 33], errors='coerce') - 1) * 100)}%"

                # Grey bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167
                # Red bar: Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167
                # Yellow Highlight: calc*=  (File “Airlines Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D167 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AF8) -1
                global_brand1 = ceil(pd.to_numeric(Airline_df.iloc[166, 3], errors='coerce')*100) / 100
                global_brand2 = ceil(pd.to_numeric(Airline_df.iloc[166, 2], errors='coerce')*100) / 100
                bencmark5 = f"{ceil((pd.to_numeric(Airline_df.iloc[166, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 31], errors='coerce') - 1) * 100)}%"

                for index, shape in enumerate(slide.shapes):
                # Yellow Highlight: calc* = (Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243 / Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243) -1
                    if index == 12 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= bencmark4
                    if index == 13 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= bencmark3
                    if index == 14 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= bencmark2
                    if index == 15 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= bencmark1
                    if index == 17 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Airline_df.iloc[242, 2], errors='coerce') / pd.to_numeric(Airline_df.iloc[242, 3], errors='coerce') - 1) * 100)}"
                #Yellow Highlight: calc* = (Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224 / Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224) -1
                    if index == 18 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Airline_df.iloc[223, 2], errors='coerce') / pd.to_numeric(Airline_df.iloc[223, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205 / Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205) -1
                    if index == 20 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Airline_df.iloc[204, 2], errors='coerce') / pd.to_numeric(Airline_df.iloc[204, 3], errors='coerce') - 1) * 100)}"
                #Yellow Highlight: calc* = (Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186 / Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186) -1
                    if index == 22 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Airline_df.iloc[185, 2], errors='coerce') / pd.to_numeric(Airline_df.iloc[185, 3], errors='coerce') - 1) * 100)}"
                #Yellow Highlight: calc* = (Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167 / Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167) -1
                    if index == 24 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Airline_df.iloc[166, 2], errors='coerce') / pd.to_numeric(Airline_df.iloc[166, 3], errors='coerce') - 1) * 100)}"
                    if index == 6 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                            if paragraph_index == 1:
                #                 Source cell: A1
                                paragraph.runs[0].text = new_data_for_source
                                for run in paragraph.runs[1:]:
                                    run.text = "" 
                            if paragraph_index == 2:
                #                 Test sample size cell: C11 Control sample size cell: D11
                                paragraph.runs[0].text = new_data_for_source1
                                for run in paragraph.runs[1:]:
                                    run.text = "" 

                chart_data = CategoryChartData()
                chart_data.categories = ['Is a market leader', 'Connects the world','Is a brand I love to fly with', 'Delivers premium in-flight experience','Is a global brand']  
                chart_data.add_series('Series 2', (market_leader2, connect_world2, brand_love2,flight_exp2,global_brand2)) 
                chart_data.add_series('Series 1', (market_leader1, connect_world1, brand_love1,flight_exp1,global_brand1)) 
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED:
                        chart = shape.chart
                        chart.replace_data(chart_data)
                        # Apply formatting to each series in the chart
                        for series in chart.series:
                            series.data_labels.number_format = '0%'  # Format as percentage
                            series.data_labels.show_value = True     # Ensure values are shown
                        break
                start_slide_index = 13
                end_slide_index = 19

                # Delete slides in reverse order
                for i in range(end_slide_index, start_slide_index - 1, -1):
                    # Remove slide by index
                    xml_slides = prs.slides._sldIdLst  # Access internal list of slide IDs
                    slides = list(xml_slides)
                    prs.part.drop_rel(slides[i].rId)  # Drop the relationship
                    del xml_slides[i]  # Delete the slide
                print("----- Slide No 13 Completed")
                prs.save('AdScore Reporting Template.pptx')
        except Exception as e:
            print(f"An error occurred at slide 13:  {e}") 

    if selected_option == "Travel":
        # # Slide 14 Travel
        # In[15]:
        # Travel Data for slide 14
        # Travel_excel = pd.ExcelFile('Travel example.xlsx')
        # Travel_csv= "Travel.csv"
        # Travel_df = pd.read_excel(Travel_excel, 'CAMP15_0902_v5', header=None)
        # Travel_df.to_csv(Travel_csv, index=False, header=False)
        # Travel_df = pd.read_csv(Travel_csv, header=None)
        Travel_df = df 
        # In[16]:
        try:
            slide_index = 13
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                new_data_for_source = Travel_df.iloc[0, 0]
                new_data_for_source =f'Source: {new_data_for_source}'
                test_value =  pd.to_numeric(Travel_df.iloc[10, 2], errors='coerce')  #C11
                control_value = pd.to_numeric(Travel_df.iloc[10, 3], errors='coerce') # D11
                new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'


                # Grey bar: Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243
                # Red bar: Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243
                # Yellow Highlight: calc*=  (File “Travel Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D243 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AX8) -1
                varietly_things1 = ceil(pd.to_numeric(Travel_df.iloc[242, 3], errors='coerce')*100) / 100
                varietly_things2 = ceil(pd.to_numeric(Travel_df.iloc[242, 2], errors='coerce')*100) / 100
                variety_bencmark1= f"{ceil((pd.to_numeric(Travel_df.iloc[242, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 49], errors='coerce') - 1) * 100)}%"

                # Grey bar: Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224
                # Red bar: Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224
                # Yellow Highlight: calc*=  (File “Travel Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D224 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AV8) -1
                vibrant_Exciting1 = ceil(pd.to_numeric(Travel_df.iloc[223, 3], errors='coerce')*100) / 100
                vibrant_Exciting2 = ceil(pd.to_numeric(Travel_df.iloc[223, 2], errors='coerce')*100) / 100
                vibrant_bencmark2 = f"{ceil((pd.to_numeric(Travel_df.iloc[223, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,47], errors='coerce') - 1) * 100)}%"

                # Grey bar: Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205
                # Red bar: Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205
                # Yellow Highlight: calc*=  (File “Travel Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D205 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AT8) -1
                safe_visit1 = ceil(pd.to_numeric(Travel_df.iloc[204, 3], errors='coerce')*100) / 100
                safe_visit2 = ceil(pd.to_numeric(Travel_df.iloc[204, 2], errors='coerce')*100) / 100
                safe_bencmark3 = f"{ceil((pd.to_numeric(Travel_df.iloc[204, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,45], errors='coerce') - 1) * 100)}%"

                # Grey bar: Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186
                # Red bar: Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186
                # Yellow Highlight: calc*=  (File “Travel Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D186 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AR8) -1

                beautiful_countryside1 = ceil(pd.to_numeric(Travel_df.iloc[185, 3], errors='coerce')*100) / 100
                beautiful_countryside2 = ceil(pd.to_numeric(Travel_df.iloc[185, 2], errors='coerce')*100) / 100
                beautiful_bencmark4 = f"{ceil((pd.to_numeric(Travel_df.iloc[185, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 43], errors='coerce') - 1) * 100)}%"

                # Grey bar: Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167
                # Red bar: Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167
                # Yellow Highlight: calc*=  (File “Travel Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D167 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AP8) -1

                cultural_heritage1 = ceil(pd.to_numeric(Travel_df.iloc[166, 3], errors='coerce')*100) / 100
                cultural_heritage2 = ceil(pd.to_numeric(Travel_df.iloc[166, 2], errors='coerce')*100) / 100
                cultural_bencmark5 = f"{ceil((pd.to_numeric(Travel_df.iloc[166, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 41], errors='coerce') - 1) * 100)}%"

                for index, shape in enumerate(slide.shapes):
                # Yellow Highlight: calc* = (Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243 / Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243) -1
                    if index == 12 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= beautiful_bencmark4
                    if index == 13 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= safe_bencmark3
                    if index == 14 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= vibrant_bencmark2
                    if index == 15 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= variety_bencmark1

                # Yellow Highlight: calc* = (Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243 / Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243) -1
                    if index == 17 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Travel_df.iloc[242, 2], errors='coerce') / pd.to_numeric(Travel_df.iloc[242, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224 / Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224) -1    
                    if index == 18 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Travel_df.iloc[223, 2], errors='coerce') / pd.to_numeric(Travel_df.iloc[223, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205 / Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205) -1
                    if index == 20 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Travel_df.iloc[204, 2], errors='coerce') / pd.to_numeric(Travel_df.iloc[204, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186 / Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186) -1
                    if index == 22 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Travel_df.iloc[185, 2], errors='coerce') / pd.to_numeric(Travel_df.iloc[185, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167 / Travel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167) -1
                    if index == 24 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Travel_df.iloc[166, 2], errors='coerce') / pd.to_numeric(Travel_df.iloc[166, 3], errors='coerce') - 1) * 100)}"
                    if index == 6 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                            if paragraph_index == 1:
                #                 Source cell: A1
                                paragraph.runs[0].text = new_data_for_source
                                for run in paragraph.runs[1:]:
                                    run.text = "" 
                            if paragraph_index == 2:
                #                 Test sample size cell: C11 Control sample size cell: D11
                                paragraph.runs[0].text = new_data_for_source1
                                for run in paragraph.runs[1:]:
                                    run.text = "" 


                chart_data = CategoryChartData()
                chart_data.categories = ['Has a variety of things to see and.', 'Is vibrant and exciting','Is safe to visit', 'Has beautiful countryside and','Has cultural heritage']  
                chart_data.add_series('Series 2', (varietly_things2, vibrant_Exciting2, safe_visit2,beautiful_countryside2,cultural_heritage2)) 
                chart_data.add_series('Series 1', (varietly_things1, vibrant_Exciting1, safe_visit1,beautiful_countryside1,cultural_heritage1)) 
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED:
                        chart = shape.chart
                        chart.replace_data(chart_data)

                        # Apply formatting to each series in the chart
                        for series in chart.series:
                            series.data_labels.number_format = '0%'  # Format as percentage
                            series.data_labels.show_value = True     # Ensure values are shown

                        break
                slides_to_delete = [12, 14, 15, 16, 17, 18, 19]

                # Sort the list in reverse order to avoid indexing issues while deleting
                slides_to_delete.sort(reverse=True)

                # Delete slides
                for slide_index in slides_to_delete:
                    xml_slides = prs.slides._sldIdLst  # internal list of slide IDs
                    slides = list(xml_slides)
                    prs.part.drop_rel(slides[slide_index].rId)  # Drop the relationship
                    del xml_slides[slide_index] 
                prs.save('AdScore Reporting Template.pptx')
        except Exception as e:
            print(f"An error occurred at slide 14:  {e}") 

    if selected_option =="Finance":
        # # slide 15 Finance

        # In[17]:


        # Finanace Data for slide 15
        # Finance_excel = pd.ExcelFile('Finance example.xlsx')
        # Finance_csv= "Finance.csv"
        # finance_df = pd.read_excel(Finance_excel, 'CAMP15_0902_v5', header=None)
        # finance_df.to_csv(Finance_csv, index=False, header=False)
        # finance_df = pd.read_csv(Finance_csv, header=None)

        finance_df = df 
        # In[18]:


        try:
            slide_index = 14  
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                new_data_for_source = finance_df.iloc[0, 0]
                new_data_for_source =f'Source: {new_data_for_source}'
                test_value =  pd.to_numeric(finance_df.iloc[10, 2], errors='coerce')  #C11
                control_value = pd.to_numeric(finance_df.iloc[10, 3], errors='coerce') # D11
                new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'

                # Grey bar: Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243
                #                 Red bar: Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243
                #                                 Yellow Highlight: calc*=  (File “Finance Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D243 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BH8) -1
                leading_instituation1 = ceil(pd.to_numeric(finance_df.iloc[242, 3], errors='coerce')*100) / 100
                leading_instituation2 = ceil(pd.to_numeric(finance_df.iloc[242, 2], errors='coerce')*100) / 100
                leading_bencmark1= f"{ceil((pd.to_numeric(finance_df.iloc[242, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 59], errors='coerce') - 1) * 100)}%"

                # Grey bar: Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224
                #                 Red bar: Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224
                #                                 Yellow Highlight: calc*=  (File “Finance Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D224 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BF8) -1
                trustworthy1 = ceil(pd.to_numeric(finance_df.iloc[223, 3], errors='coerce')*100) / 100
                trustworthy2 = ceil(pd.to_numeric(finance_df.iloc[223, 2], errors='coerce')*100) / 100
                trustworthy_bencmark2 = f"{ceil((pd.to_numeric(finance_df.iloc[223, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,57], errors='coerce') - 1) * 100)}%"

                # Grey bar: Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205
                #                 Red bar: Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205
                #                                 Yellow Highlight: calc*=  (File “Finance Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D205 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BD8) -1
                customer_focus1 = ceil(pd.to_numeric(finance_df.iloc[204, 3], errors='coerce')*100) / 100
                customer_focus2 = ceil(pd.to_numeric(finance_df.iloc[204, 2], errors='coerce')*100) / 100
                customer_bencmark3 = f"{ceil((pd.to_numeric(finance_df.iloc[204, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,55], errors='coerce') - 1) * 100)}%"
                # Grey bar: Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186
                #                 Red bar: Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186
                #                                 Yellow Highlight: calc*=  (File “Finance Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D186 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BB8) -1

                global_brand1 = ceil(pd.to_numeric(finance_df.iloc[185, 3], errors='coerce')*100) / 100
                global_brand2 = ceil(pd.to_numeric(finance_df.iloc[185, 2], errors='coerce')*100) / 100
                global_bencmark4 = f"{ceil((pd.to_numeric(finance_df.iloc[185, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 53], errors='coerce') - 1) * 100)}%"

                # Grey bar: Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167
                #                 Red bar: Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167
                #                                 Yellow Highlight: calc*=  (File “Finance Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D167 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: AZ8) -1
                service_provide1 = ceil(pd.to_numeric(finance_df.iloc[166, 3], errors='coerce')*100) / 100
                service_provide2 = ceil(pd.to_numeric(finance_df.iloc[166, 2], errors='coerce')*100) / 100
                service_bencmark5 = f"{ceil((pd.to_numeric(finance_df.iloc[166, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 51], errors='coerce') - 1) * 100)}%"

                for index, shape in enumerate(slide.shapes):
                # Yellow Highlight: calc* = (Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243 / Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243) -1
                    if index == 12 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= global_bencmark4
                    if index == 13 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= customer_bencmark3
                    if index == 14 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= trustworthy_bencmark2
                    if index == 15 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= leading_bencmark1

                # Yellow Highlight: calc* = (Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243 / Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243) -1
                    if index == 17 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(finance_df.iloc[242, 2], errors='coerce') / pd.to_numeric(finance_df.iloc[242, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224 / Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224) -1
                    if index == 18 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(finance_df.iloc[223, 2], errors='coerce') / pd.to_numeric(finance_df.iloc[223, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205 / Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205) -1
                    if index == 20 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(finance_df.iloc[204, 2], errors='coerce') / pd.to_numeric(finance_df.iloc[204, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186 / Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186) -1
                    if index == 22 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(finance_df.iloc[185, 2], errors='coerce') / pd.to_numeric(finance_df.iloc[185, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167 / Finance Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167) -1
                    if index == 24 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(finance_df.iloc[166, 2], errors='coerce') / pd.to_numeric(finance_df.iloc[166, 3], errors='coerce') - 1) * 100)}"
                    if index == 6 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                            if paragraph_index == 1:
                #                 Source cell: A1
                                paragraph.runs[0].text = new_data_for_source
                                for run in paragraph.runs[1:]:
                                    run.text = "" 
                            if paragraph_index == 2:
                #                 Test sample size cell: C11 Control sample size cell: D11
                                paragraph.runs[0].text = new_data_for_source1
                                for run in paragraph.runs[1:]:
                                    run.text = "" 


                chart_data = CategoryChartData()
                chart_data.categories = ['Has a variety of things to see and.', 'Is vibrant and exciting','Is safe to visit', 'Has beautiful countryside and','Has cultural heritage']  
                chart_data.add_series('Series 2', (leading_instituation2, trustworthy2, customer_focus2,global_brand2,service_provide2)) 
                chart_data.add_series('Series 1', (leading_instituation1, trustworthy1, customer_focus1,global_brand1,service_provide1)) 
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED:
                        chart = shape.chart
                        chart.replace_data(chart_data)

                        # Apply formatting to each series in the chart
                        for series in chart.series:
                            series.data_labels.number_format = '0%'  # Format as percentage
                            series.data_labels.show_value = True     # Ensure values are shown

                        break
                slides_to_delete = [12, 13, 15, 16, 17, 18, 19]

                # Sort the list in reverse order to avoid indexing issues while deleting
                slides_to_delete.sort(reverse=True)

                # Delete slides
                for slide_index in slides_to_delete:
                    xml_slides = prs.slides._sldIdLst  # internal list of slide IDs
                    slides = list(xml_slides)
                    prs.part.drop_rel(slides[slide_index].rId)  # Drop the relationship
                    del xml_slides[slide_index] 

                prs.save('AdScore Reporting Template.pptx')
        except Exception as e:
            print(f"An error occurred at slide 15:  {e}") 

    if selected_option == "Hotel":
        # # Slide 16 Hotel 
        # In[19]:
        # Hotel Data for slide 16
        # Hotel_excel = pd.ExcelFile('Hotel example.xlsx')
        # Hotel_csv= "Hotel.csv"
        # Hotel_df = pd.read_excel(Hotel_excel, 'CAMP15_0902_v5', header=None)
        # Hotel_df.to_csv(Hotel_csv, index=False, header=False)
        # Hotel_df = pd.read_csv(Hotel_csv, header=None)

        Hotel_df  = df

        # In[20]:
        try:
            slide_index = 15
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                new_data_for_source = Hotel_df.iloc[0, 0]
                new_data_for_source =f'Source: {new_data_for_source}'
                test_value =  pd.to_numeric(Hotel_df.iloc[10, 2], errors='coerce')  #C11
                control_value = pd.to_numeric(Hotel_df.iloc[10, 3], errors='coerce') # D11
                new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'

                # Grey bar: Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243
                #                 Red bar: Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243
                #                                 Yellow Highlight: calc*=  (File “Hotel Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D243 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BP8) -1
                great_location1 = ceil(pd.to_numeric(Hotel_df.iloc[242, 3], errors='coerce')*100) / 100
                great_location2 = ceil(pd.to_numeric(Hotel_df.iloc[242, 2], errors='coerce')*100) / 100
                location_bencmark1= f"{ceil((pd.to_numeric(Hotel_df.iloc[242, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 67], errors='coerce') - 1) * 100)}%"

                # Grey bar: Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224
                #                 Red bar: Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224
                #                                 Yellow Highlight: calc*=  (File “Hotel Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D224 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BP8) -1
                luxurios1 = ceil(pd.to_numeric(Hotel_df.iloc[223, 3], errors='coerce')*100) / 100
                luxurios2 = ceil(pd.to_numeric(Hotel_df.iloc[223, 2], errors='coerce')*100) / 100
                luxurios_bencmark2 = f"{ceil((pd.to_numeric(Hotel_df.iloc[223, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,67], errors='coerce') - 1) * 100)}%"
                # Grey bar: Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205
                #                 Red bar: Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205
                #                                 Yellow Highlight: calc*=  (File “Hotel Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D205 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BN8) -1
                love_place1 = ceil(pd.to_numeric(Hotel_df.iloc[204, 3], errors='coerce')*100) / 100
                love_place2 = ceil(pd.to_numeric(Hotel_df.iloc[204, 2], errors='coerce')*100) / 100
                love_place_bencmark3 = f"{ceil((pd.to_numeric(Hotel_df.iloc[204, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,65], errors='coerce') - 1) * 100)}%"
                # Grey bar: Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186
                #                 Red bar: Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186
                #                                 Yellow Highlight: calc*=  (File “Hotel Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D186 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BL8) -1

                experience1 = ceil(pd.to_numeric(Hotel_df.iloc[185, 3], errors='coerce')*100) / 100
                experience2 = ceil(pd.to_numeric(Hotel_df.iloc[185, 2], errors='coerce')*100) / 100
                experience_bencmark4 = f"{ceil((pd.to_numeric(Hotel_df.iloc[185, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 63], errors='coerce') - 1) * 100)}%"

                # Grey bar: Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167
                #                 Red bar: Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167
                #                                 Yellow Highlight: calc*=  (File “Hotel Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D167 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BJ8) -1
                global_brand1 = ceil(pd.to_numeric(Hotel_df.iloc[166, 3], errors='coerce')*100) / 100
                global_brand2 = ceil(pd.to_numeric(Hotel_df.iloc[166, 2], errors='coerce')*100) / 100
                global_brand_bencmark5 = f"{ceil((pd.to_numeric(Hotel_df.iloc[166, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 61], errors='coerce') - 1) * 100)}%"

                for index, shape in enumerate(slide.shapes):
                # Yellow Highlight: calc* = (Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243 / Airlines Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243) -1
                    if index == 12 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= experience_bencmark4
                    if index == 13 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= love_place_bencmark3
                    if index == 14 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= luxurios_bencmark2
                    if index == 15 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= location_bencmark1

                # Yellow Highlight: calc* = (Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243 / Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243) -1
                    if index == 17 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Hotel_df.iloc[242, 2], errors='coerce') / pd.to_numeric(Hotel_df.iloc[242, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224 / Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224) -1
                    if index == 18 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Hotel_df.iloc[223, 2], errors='coerce') / pd.to_numeric(Hotel_df.iloc[223, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205 / Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205) -1
                    if index == 20 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Hotel_df.iloc[204, 2], errors='coerce') / pd.to_numeric(Hotel_df.iloc[204, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186 / Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186) -1
                    if index == 22 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Hotel_df.iloc[185, 2], errors='coerce') / pd.to_numeric(Hotel_df.iloc[185, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167 / Hotel Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167) -1
                    if index == 24 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(Hotel_df.iloc[166, 2], errors='coerce') / pd.to_numeric(Hotel_df.iloc[166, 3], errors='coerce') - 1) * 100)}"
                    if index == 6 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                            if paragraph_index == 1:
                #                 Source cell: A1
                                paragraph.runs[0].text = new_data_for_source
                                for run in paragraph.runs[1:]:
                                    run.text = "" 
                            if paragraph_index == 2:
                #                 Test sample size cell: C11 Control sample size cell: D11
                                paragraph.runs[0].text = new_data_for_source1
                                for run in paragraph.runs[1:]:
                                    run.text = "" 


                chart_data = CategoryChartData()
                chart_data.categories = ['Has a variety of things to see and.', 'Is vibrant and exciting','Is safe to visit', 'Has beautiful countryside and','Has cultural heritage']  
                chart_data.add_series('Series 2', (great_location2, luxurios2, love_place2,experience2,global_brand2)) 
                chart_data.add_series('Series 1', (great_location1, luxurios1, love_place1,experience1,global_brand1)) 
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED:
                        chart = shape.chart
                        chart.replace_data(chart_data)

                        # Apply formatting to each series in the chart
                        for series in chart.series:
                            series.data_labels.number_format = '0%'  # Format as percentage
                            series.data_labels.show_value = True     # Ensure values are shown

                        break
                slides_to_delete = [12, 13,14, 16, 17, 18, 19]
                # Sort the list in reverse order to avoid indexing issues while deleting
                slides_to_delete.sort(reverse=True)
                # Delete slides
                for slide_index in slides_to_delete:
                    xml_slides = prs.slides._sldIdLst  # internal list of slide IDs
                    slides = list(xml_slides)
                    prs.part.drop_rel(slides[slide_index].rId)  # Drop the relationship
                    del xml_slides[slide_index] 
                prs.save('AdScore Reporting Template.pptx')
        except Exception as e:
            print(f"An error occurred at slide 16:  {e}") 

    if selected_option == "Auto":
        # # Slide 17 Auto

        # In[21]:


        # auto Data for slide 16
        # auto_excel = pd.ExcelFile('Automotive example.xlsx')
        # auto_csv= "auto.csv"
        # auto_df = pd.read_excel(auto_excel, 'CAMP15_0902_v5', header=None)
        # auto_df.to_csv(auto_csv, index=False, header=False)
        # auto_df = pd.read_csv(auto_csv, header=None)

        auto_df = df
        # In[22]:


        try:
            slide_index = 16  
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                new_data_for_source = auto_df.iloc[0, 0]
                new_data_for_source =f'Source: {new_data_for_source}'
                test_value =  pd.to_numeric(auto_df.iloc[10, 2], errors='coerce')  #C11
                control_value = pd.to_numeric(auto_df.iloc[10, 3], errors='coerce') # D11
                new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'

                # Grey bar: Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243
                #                 Red bar: Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243
                #                                 Yellow Highlight: calc*=  (File “Automotive Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D243 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CB6) -1
                vechicle_peoplelike1 = ceil(pd.to_numeric(auto_df.iloc[242, 3], errors='coerce')*100) / 100
                vechicle_peoplelike2 = ceil(pd.to_numeric(auto_df.iloc[242, 2], errors='coerce')*100) / 100
                vechicle_peoplelike_bencmark1= f"{ceil((pd.to_numeric(auto_df.iloc[242, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 79], errors='coerce') - 1) * 100)}%"

                # Grey bar: Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224
                #                 Red bar: Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224
                #                                 Yellow Highlight: calc*=  (File “Automotive Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D224 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BZ6) -1
                vechicle_enjoyable1 = ceil(pd.to_numeric(auto_df.iloc[223, 3], errors='coerce')*100) / 100
                vechicle_enjoyable2 = ceil(pd.to_numeric(auto_df.iloc[223, 2], errors='coerce')*100) / 100
                vechicle_enjoyable_bencmark2 = f"{ceil((pd.to_numeric(auto_df.iloc[223, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,77], errors='coerce') - 1) * 100)}%"
                # Grey bar: Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205
                #                 Red bar: Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205
                #                                 Yellow Highlight: calc*=  (File “Automotive Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D205 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BX6) -1
                proud_vechicle1 = ceil(pd.to_numeric(auto_df.iloc[204, 3], errors='coerce')*100) / 100
                proud_vechicle2 = ceil(pd.to_numeric(auto_df.iloc[204, 2], errors='coerce')*100) / 100
                proud_vechicle_bencmark3 = f"{ceil((pd.to_numeric(auto_df.iloc[204, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,75], errors='coerce') - 1) * 100)}%"
                # Grey bar: Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186
                #                 Red bar: Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186
                #                                 Yellow Highlight: calc*=  (File “Automotive Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D186 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BV6) -1

                high_quality1 = ceil(pd.to_numeric(auto_df.iloc[185, 3], errors='coerce')*100) / 100
                high_quality2 = ceil(pd.to_numeric(auto_df.iloc[185, 2], errors='coerce')*100) / 100
                high_quality_bencmark4 = f"{ceil((pd.to_numeric(auto_df.iloc[185, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 73], errors='coerce') - 1) * 100)}%"
                # Grey bar: Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167
                #                 Red bar: Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167
                #                                 Yellow Highlight: calc*=  (File “Automotive Automotive.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D167 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: BT6) -1

                brand1 = ceil(pd.to_numeric(auto_df.iloc[166, 3], errors='coerce')*100) / 100
                brand2 = ceil(pd.to_numeric(auto_df.iloc[166, 2], errors='coerce')*100) / 100
                brand_bencmark5 = f"{ceil((pd.to_numeric(auto_df.iloc[166, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 71], errors='coerce') - 1) * 100)}%"

                for index, shape in enumerate(slide.shapes):
                    if index == 12 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= high_quality_bencmark4
                    if index == 13 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= proud_vechicle_bencmark3
                    if index == 14 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= vechicle_enjoyable_bencmark2
                    if index == 15 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= vechicle_peoplelike_bencmark1

                # Yellow Highlight: calc* = (Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243 / Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243) -1
                    if index == 17 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(auto_df.iloc[242, 2], errors='coerce') / pd.to_numeric(auto_df.iloc[242, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: alc* = (Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224 / Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224) -1
                    if index == 18 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(auto_df.iloc[223, 2], errors='coerce') / pd.to_numeric(auto_df.iloc[223, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205 / Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205) -1
                    if index == 20 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(auto_df.iloc[204, 2], errors='coerce') / pd.to_numeric(auto_df.iloc[204, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186 / Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186) -1
                    if index == 22 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(auto_df.iloc[185, 2], errors='coerce') / pd.to_numeric(auto_df.iloc[185, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167 / Automotive Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167) -1
                    if index == 24 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(auto_df.iloc[166, 2], errors='coerce') / pd.to_numeric(auto_df.iloc[166, 3], errors='coerce') - 1) * 100)}"
                    if index == 6 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                            if paragraph_index == 1:
                #                 Source cell: A1
                                paragraph.runs[0].text = new_data_for_source
                                for run in paragraph.runs[1:]:
                                    run.text = "" 
                            if paragraph_index == 2:
                #                 Test sample size cell: C11 Control sample size cell: D11
                                paragraph.runs[0].text = new_data_for_source1
                                for run in paragraph.runs[1:]:
                                    run.text = "" 


                chart_data = CategoryChartData()
                chart_data.categories = ['Has a variety of things to see and.', 'Is vibrant and exciting','Is safe to visit', 'Has beautiful countryside and','Has cultural heritage']  
                chart_data.add_series('Series 2', (vechicle_peoplelike2, vechicle_enjoyable2, proud_vechicle2,high_quality2,brand2)) 
                chart_data.add_series('Series 1', (vechicle_peoplelike1, vechicle_enjoyable1, proud_vechicle1,high_quality1,brand1)) 
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED:
                        chart = shape.chart
                        chart.replace_data(chart_data)

                        # Apply formatting to each series in the chart
                        for series in chart.series:
                            series.data_labels.number_format = '0%'  # Format as percentage
                            series.data_labels.show_value = True     # Ensure values are shown

                        break
                slides_to_delete = [12, 13,14, 15, 17, 18, 19]

                # Sort the list in reverse order to avoid indexing issues while deleting
                slides_to_delete.sort(reverse=True)

                # Delete slides
                for slide_index in slides_to_delete:
                    xml_slides = prs.slides._sldIdLst  # internal list of slide IDs
                    slides = list(xml_slides)
                    prs.part.drop_rel(slides[slide_index].rId)  # Drop the relationship
                    del xml_slides[slide_index] 
                prs.save('AdScore Reporting Template.pptx')
        except Exception as e:
            print(f"An error occurred at slide 17:  {e}") 

    if selected_option =="Tech":
        # # Slide 18 Tech

        # In[23]:


        # tech Data for slide 18
        # tech_excel = pd.ExcelFile('Tech example.xlsx')
        # tech_csv= "tech.csv"
        # tech_df = pd.read_excel(tech_excel, 'CAMP15_0902_v5', header=None)
        # tech_df.to_csv(tech_csv, index=False, header=False)
        # tech_df = pd.read_csv(tech_csv, header=None)

        tech_df = df
        # In[24]:


        try:
            # Check if the slide index exists
            slide_index = 17  # Example slide index
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                new_data_for_source = tech_df.iloc[0, 0]
                new_data_for_source =f'Source: {new_data_for_source}'
                test_value =  pd.to_numeric(tech_df.iloc[10, 2], errors='coerce')  #C11
                control_value = pd.to_numeric(tech_df.iloc[10, 3], errors='coerce') # D11
                new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'

                # Grey bar: Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243
                #                 Red bar: Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243
                #                                 Yellow Highlight: calc*=  (File “Tech Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D243 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CL8) -1
                premium_brand1 = ceil(pd.to_numeric(tech_df.iloc[242, 3], errors='coerce')*100) / 100
                premium_brand2 = ceil(pd.to_numeric(tech_df.iloc[242, 2], errors='coerce')*100) / 100
                premium_brand_bencmark1= f"{ceil((pd.to_numeric(tech_df.iloc[242, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 89], errors='coerce') - 1) * 100)}%"

                # Grey bar: Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224
                #                 Red bar: Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224
                #                                 Yellow Highlight: calc*=  (File “Tech Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D224 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CJ8) -1
                quality_tech1 = ceil(pd.to_numeric(tech_df.iloc[223, 3], errors='coerce')*100) / 100
                quality_tech2 = ceil(pd.to_numeric(tech_df.iloc[223, 2], errors='coerce')*100) / 100
                quality_tech_bencmark2 = f"{ceil((pd.to_numeric(tech_df.iloc[223, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,87], errors='coerce') - 1) * 100)}%"
                # Grey bar: Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205
                #                 Red bar: Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205
                #                                 Yellow Highlight: calc*=  (File “Tech Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D205 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CH8) -1
                innovative_tech1 = ceil(pd.to_numeric(tech_df.iloc[204, 3], errors='coerce')*100) / 100
                innovative_tech2 = ceil(pd.to_numeric(tech_df.iloc[204, 2], errors='coerce')*100) / 100
                innovative_tech_bencmark3 = f"{ceil((pd.to_numeric(tech_df.iloc[204, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,85], errors='coerce') - 1) * 100)}%"
                # Grey bar: Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186
                #                 Red bar: Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186
                #                                 Yellow Highlight: calc*=  (File “Tech Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D186 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CF8) -1

                user_exp1 = ceil(pd.to_numeric(tech_df.iloc[185, 3], errors='coerce')*100) / 100
                user_exp2 = ceil(pd.to_numeric(tech_df.iloc[185, 2], errors='coerce')*100) / 100
                user_exp_bencmark4 = f"{ceil((pd.to_numeric(tech_df.iloc[185, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 83], errors='coerce') - 1) * 100)}%"
                # Grey bar: Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167
                #                 Red bar: Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167
                #                                 Yellow Highlight: calc*=  (File “Tech Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D167 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CD8) -1

                people_tech1 = ceil(pd.to_numeric(tech_df.iloc[166, 3], errors='coerce')*100) / 100
                people_tech2 = ceil(pd.to_numeric(tech_df.iloc[166, 2], errors='coerce')*100) / 100
                people_tech_bencmark5 = f"{ceil((pd.to_numeric(tech_df.iloc[166, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 81], errors='coerce') - 1) * 100)}%"

                for index, shape in enumerate(slide.shapes):
                    if index == 12 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= user_exp_bencmark4
                    if index == 13 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= innovative_tech_bencmark3
                    if index == 14 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= quality_tech_bencmark2
                    if index == 15 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= premium_brand_bencmark1

                # Yellow Highlight: calc* = (Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243 / Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243) -1
                    if index == 17 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(tech_df.iloc[242, 2], errors='coerce') / pd.to_numeric(tech_df.iloc[242, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224 / Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224) -1
                    if index == 18 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(tech_df.iloc[223, 2], errors='coerce') / pd.to_numeric(tech_df.iloc[223, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205 / Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205) -1
                    if index == 20 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(tech_df.iloc[204, 2], errors='coerce') / pd.to_numeric(tech_df.iloc[204, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186 / Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186) -1
                    if index == 22 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(tech_df.iloc[185, 2], errors='coerce') / pd.to_numeric(tech_df.iloc[185, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167 / Tech Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167) -1
                    if index == 24 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(tech_df.iloc[166, 2], errors='coerce') / pd.to_numeric(tech_df.iloc[166, 3], errors='coerce') - 1) * 100)}"
                    if index == 6 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                            if paragraph_index == 1:
                #                 Source cell: A1
                                paragraph.runs[0].text = new_data_for_source
                                for run in paragraph.runs[1:]:
                                    run.text = "" 
                            if paragraph_index == 2:
                #                 Test sample size cell: C11 Control sample size cell: D11
                                paragraph.runs[0].text = new_data_for_source1
                                for run in paragraph.runs[1:]:
                                    run.text = "" 


                chart_data = CategoryChartData()
                chart_data.categories = ['Has a variety of things to see and.', 'Is vibrant and exciting','Is safe to visit', 'Has beautiful countryside and','Has cultural heritage']  
                chart_data.add_series('Series 2', (premium_brand2, quality_tech2, innovative_tech2,user_exp2,people_tech2)) 
                chart_data.add_series('Series 1', (premium_brand1, quality_tech1, innovative_tech1,user_exp1,people_tech1)) 
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED:
                        chart = shape.chart
                        chart.replace_data(chart_data)

                        # Apply formatting to each series in the chart
                        for series in chart.series:
                            series.data_labels.number_format = '0%'  # Format as percentage
                            series.data_labels.show_value = True     # Ensure values are shown

                        break
                slides_to_delete = [12, 13,14, 15, 16, 18, 19]
                # Sort the list in reverse order to avoid indexing issues while deleting
                slides_to_delete.sort(reverse=True)

                # Delete slides
                for slide_index in slides_to_delete:
                    xml_slides = prs.slides._sldIdLst  # internal list of slide IDs
                    slides = list(xml_slides)
                    prs.part.drop_rel(slides[slide_index].rId)  # Drop the relationship
                    del xml_slides[slide_index] 
                prs.save('AdScore Reporting Template.pptx')
        except Exception as e:
            print(f"An error occurred at slide 18:  {e}") 

    if selected_option =="FDI":
        # # Slide 19 FDI

        # In[25]:


        # FDI Data for slide 19
        # fdi_excel = pd.ExcelFile('Foreign Trade & Investment example.xlsx')
        # fdi_csv= "FDI.csv"
        # fdi_df = pd.read_excel(fdi_excel, 'CAMP15_0902_v5', header=None)
        # fdi_df.to_csv(fdi_csv, index=False, header=False)
        # fdi_df = pd.read_csv(fdi_csv, header=None)

        fdi_df = df
        # In[26]:


        try:
            # Check if the slide index exists
            slide_index = 18  # Example slide index
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                new_data_for_source = fdi_df.iloc[0, 0]
                new_data_for_source =f'Source: {new_data_for_source}'
                test_value =  pd.to_numeric(fdi_df.iloc[10, 2], errors='coerce')  #C11
                control_value = pd.to_numeric(fdi_df.iloc[10, 3], errors='coerce') # D11
                new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'

                # Grey bar: Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243
                #                 Red bar: Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243
                #                                 Yellow Highlight: calc*=  (File “Foreign trade & Investment Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D243 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CV6) -1
                # CV6 has no data i am assuming it is CV5
                fdi_economy1 = ceil(pd.to_numeric(fdi_df.iloc[242, 3], errors='coerce')*100) / 100
                fdi_economy2 = ceil(pd.to_numeric(fdi_df.iloc[242, 2], errors='coerce')*100) / 100
                fdi_economy_bencmark1= f"{ceil((pd.to_numeric(fdi_df.iloc[242, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[4, 99], errors='coerce') - 1) * 100)}%"

                # Grey bar: Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224
                #                 Red bar: Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224
                #                                 Yellow Highlight: calc*=  (File “Foreign trade & Investment Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D224 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CT6) -1
                fdi_place1 = ceil(pd.to_numeric(fdi_df.iloc[223, 3], errors='coerce')*100) / 100
                fdi_place2 = ceil(pd.to_numeric(fdi_df.iloc[223, 2], errors='coerce')*100) / 100
                fdi_place_bencmark2 = f"{ceil((pd.to_numeric(fdi_df.iloc[223, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[4,97], errors='coerce') - 1) * 100)}%"

                # Grey bar: Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205
                #                 Red bar: Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205
                #                                 Yellow Highlight: calc*=  (File “Foreign trade & Investment Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D205 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CR6) -1
                fdi_workforce1 = ceil(pd.to_numeric(fdi_df.iloc[204, 3], errors='coerce')*100) / 100
                fdi_workforce2 = ceil(pd.to_numeric(fdi_df.iloc[204, 2], errors='coerce')*100) / 100
                fdi_workforce_tech_bencmark3 = f"{ceil((pd.to_numeric(fdi_df.iloc[204, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[4,95], errors='coerce') - 1) * 100)}%"
                # Grey bar: Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186
                #                 Red bar: Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186
                #                                 Yellow Highlight: calc*=  (File “Foreign trade & Investment Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D186 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CP6) -1

                fdi_world1 = ceil(pd.to_numeric(fdi_df.iloc[185, 3], errors='coerce')*100) / 100
                fdi_world2 = ceil(pd.to_numeric(fdi_df.iloc[185, 2], errors='coerce')*100) / 100
                fdi_world_bencmark4 = f"{ceil((pd.to_numeric(fdi_df.iloc[185, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[4, 93], errors='coerce') - 1) * 100)}%"
                # Grey bar: Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167
                #                 Red bar: Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167
                #                                 Yellow Highlight: calc*=  (File “Foreign trade & Investment Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D167 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CN6) -1

                fdi_investor1 = ceil(pd.to_numeric(fdi_df.iloc[166, 3], errors='coerce')*100) / 100
                fdi_investor2 = ceil(pd.to_numeric(fdi_df.iloc[166, 2], errors='coerce')*100) / 100
                fdi_investor_bencmark5 = f"{ceil((pd.to_numeric(fdi_df.iloc[166, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[4, 91], errors='coerce') - 1) * 100)}%"

                for index, shape in enumerate(slide.shapes):
                    if index == 12 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= fdi_world_bencmark4
                    if index == 13 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= fdi_workforce_tech_bencmark3
                    if index == 14 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= fdi_place_bencmark2
                    if index == 15 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= fdi_economy_bencmark1

                # Yellow Highlight: calc* = (Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243 / Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243) -1
                    if index == 17 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(fdi_df.iloc[242, 2], errors='coerce') / pd.to_numeric(fdi_df.iloc[242, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224 / Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224) -1
                    if index == 18 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(fdi_df.iloc[223, 2], errors='coerce') / pd.to_numeric(fdi_df.iloc[223, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205 / Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205) -1
                    if index == 20 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(fdi_df.iloc[204, 2], errors='coerce') / pd.to_numeric(fdi_df.iloc[204, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186 / Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186) -1
                    if index == 22 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(fdi_df.iloc[185, 2], errors='coerce') / pd.to_numeric(fdi_df.iloc[185, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167 / Foreign trade & Investment Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167) -1
                    if index == 24 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(fdi_df.iloc[166, 2], errors='coerce') / pd.to_numeric(fdi_df.iloc[166, 3], errors='coerce') - 1) * 100)}"
                    if index == 6 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                            if paragraph_index == 1:
                #                 Source cell: A1
                                paragraph.runs[0].text = new_data_for_source
                                for run in paragraph.runs[1:]:
                                    run.text = "" 
                            if paragraph_index == 2:
                #                 Test sample size cell: C11 Control sample size cell: D11
                                paragraph.runs[0].text = new_data_for_source1
                                for run in paragraph.runs[1:]:
                                    run.text = "" 


                chart_data = CategoryChartData()
                chart_data.categories = ['Has a variety of things to see and.', 'Is vibrant and exciting','Is safe to visit', 'Has beautiful countryside and','Has cultural heritage']  
                chart_data.add_series('Series 2', (fdi_economy2, fdi_place2, fdi_workforce2,fdi_world2,fdi_investor2)) 
                chart_data.add_series('Series 1', (fdi_economy1, fdi_place1, fdi_workforce1,fdi_world1,fdi_investor1)) 
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED:
                        chart = shape.chart
                        chart.replace_data(chart_data)

                        # Apply formatting to each series in the chart
                        for series in chart.series:
                            series.data_labels.number_format = '0%'  # Format as percentage
                            series.data_labels.show_value = True     # Ensure values are shown

                        break
                slides_to_delete = [12, 13,14, 15, 16, 17,19]

                # Sort the list in reverse order to avoid indexing issues while deleting
                slides_to_delete.sort(reverse=True)

                # Delete slides
                for slide_index in slides_to_delete:
                    xml_slides = prs.slides._sldIdLst  # internal list of slide IDs
                    slides = list(xml_slides)
                    prs.part.drop_rel(slides[slide_index].rId)  # Drop the relationship
                    del xml_slides[slide_index] 
                prs.save('AdScore Reporting Template.pptx')
        except Exception as e:
            print(f"An error occurred at slide 19:  {e}") 

    if selected_option =="Others":

        # # Slide 20 Other

        # In[27]:


        # other Data for slide 20
        # other_excel = pd.ExcelFile('Other example.xlsx')
        # other_csv= "other.csv"
        # other_df = pd.read_excel(other_excel, 'CAMP15_0902_v5', header=None)
        # other_df.to_csv(other_csv, index=False, header=False)
        # other_df = pd.read_csv(other_csv, header=None)

        other_df = df 
        # In[28]:
        try:
            # Check if the slide index exists
            slide_index = 19  # Example slide index
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]
                new_data_for_source = other_df.iloc[0, 0]
                new_data_for_source =f'Source: {new_data_for_source}'
                test_value =  pd.to_numeric(other_df.iloc[10, 2], errors='coerce')  #C11
                control_value = pd.to_numeric(other_df.iloc[10, 3], errors='coerce') # D11
                new_data_for_source1 =f'Sample Test n={test_value} , Control n={control_value} *Caution, low base size'

                # Grey bar: Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243
                #                 Red bar: Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243
                #                                 Yellow Highlight: calc*=  (File “Other Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D243 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: DF8) -1
                other_leader1 = ceil(pd.to_numeric(other_df.iloc[242, 3], errors='coerce')*100) / 100
                other_leader2 = ceil(pd.to_numeric(other_df.iloc[242, 2], errors='coerce')*100) / 100
                other_leader_bencmark1= f"{ceil((pd.to_numeric(other_df.iloc[242, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 109], errors='coerce') - 1) * 100)}%"

                # Grey bar: Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224
                #                 Red bar: Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224
                #                                 Yellow Highlight: calc*=  (File “Other Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D224 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: DD8) -1
                other_luxurious1 = ceil(pd.to_numeric(other_df.iloc[223, 3], errors='coerce')*100) / 100
                other_luxurious2 = ceil(pd.to_numeric(other_df.iloc[223, 2], errors='coerce')*100) / 100
                other_luxurious_bencmark2 = f"{ceil((pd.to_numeric(other_df.iloc[223, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,107], errors='coerce') - 1) * 100)}%"

                # Grey bar: Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205
                #                 Red bar: Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205
                #                                 Yellow Highlight: calc*=  (File “Other Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D205 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: DB8) -1
                other_love1 = ceil(pd.to_numeric(other_df.iloc[204, 3], errors='coerce')*100) / 100
                other_love2 = ceil(pd.to_numeric(other_df.iloc[204, 2], errors='coerce')*100) / 100
                other_love_bencmark3 = f"{ceil((pd.to_numeric(other_df.iloc[204, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7,105], errors='coerce') - 1) * 100)}%"
                # Grey bar: Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186
                #                 Red bar: Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186
                #                                 Yellow Highlight: calc*=  (File “Other Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D186 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CZ8) -1

                other_quality1 = ceil(pd.to_numeric(other_df.iloc[185, 3], errors='coerce')*100) / 100
                other_quality2 = ceil(pd.to_numeric(other_df.iloc[185, 2], errors='coerce')*100) / 100
                other_quality_bencmark4 = f"{ceil((pd.to_numeric(other_df.iloc[185, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 103], errors='coerce') - 1) * 100)}%"
                # Grey bar: Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167
                #                 Red bar: Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167
                #                                 Yellow Highlight: calc*=  (File “Other Example.xlsx”, Tab: “CAMP15_0902_v5”, Cell:  D167 / File “AdScore Norms W46 (W87).xlcm”, Tab: “Data”, Cell: CX8) -1

                other_brand1 = ceil(pd.to_numeric(other_df.iloc[166, 3], errors='coerce')*100) / 100
                other_brand2 = ceil(pd.to_numeric(other_df.iloc[166, 2], errors='coerce')*100) / 100
                other_brand_bencmark5 = f"{ceil((pd.to_numeric(other_df.iloc[166, 3], errors='coerce') / pd.to_numeric(AdScore_df.iloc[7, 101], errors='coerce') - 1) * 100)}%"

                for index, shape in enumerate(slide.shapes):
                    if index == 12 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= other_quality_bencmark4
                    if index == 13 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= other_love_bencmark3
                    if index == 14 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= other_luxurious_bencmark2
                    if index == 15 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= other_leader_bencmark1

                # Yellow Highlight: calc* = (Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C243 / Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D243) -1
                    if index == 17 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(other_df.iloc[242, 2], errors='coerce') / pd.to_numeric(other_df.iloc[242, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C224 / Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D224) -1
                    if index == 18 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(other_df.iloc[223, 2], errors='coerce') / pd.to_numeric(other_df.iloc[223, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C205 / Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D205) -1
                    if index == 20 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(other_df.iloc[204, 2], errors='coerce') / pd.to_numeric(other_df.iloc[204, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C186 / Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D186) -1
                    if index == 22 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(other_df.iloc[185, 2], errors='coerce') / pd.to_numeric(other_df.iloc[185, 3], errors='coerce') - 1) * 100)}"
                # Yellow Highlight: calc* = (Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  C167 / Other Example.xlsx, Tab: CAMP15_0902_v5: Cell:  D167) -1
                    if index == 24 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        paragraph.runs[0].text= f"{ceil((pd.to_numeric(other_df.iloc[166, 2], errors='coerce') / pd.to_numeric(other_df.iloc[166, 3], errors='coerce') - 1) * 100)}"
                    if index == 6 and shape.has_text_frame:
                        paragraph = shape.text_frame.paragraphs[0]
                        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                            if paragraph_index == 1:
                #                 Source cell: A1
                                paragraph.runs[0].text = new_data_for_source
                                for run in paragraph.runs[1:]:
                                    run.text = "" 
                            if paragraph_index == 2:
                #                 Test sample size cell: C11 Control sample size cell: D11
                                paragraph.runs[0].text = new_data_for_source1
                                for run in paragraph.runs[1:]:
                                    run.text = "" 


                chart_data = CategoryChartData()
                chart_data.categories = ['Has a variety of things to see and.', 'Is vibrant and exciting','Is safe to visit', 'Has beautiful countryside and','Has cultural heritage']  
                chart_data.add_series('Series 2', (other_leader2, other_luxurious2, other_love2,other_quality2,other_brand2)) 
                chart_data.add_series('Series 1', (other_leader1, other_luxurious1, other_love1,other_quality1,other_brand1)) 
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART and shape.chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED:
                        chart = shape.chart
                        chart.replace_data(chart_data)

                        # Apply formatting to each series in the chart
                        for series in chart.series:
                            series.data_labels.number_format = '0%'  # Format as percentage
                            series.data_labels.show_value = True     # Ensure values are shown

                        break
                slides_to_delete = [12, 13, 14, 15, 16, 17, 18]

                # Sort the list in reverse order to avoid indexing issues while deleting
                slides_to_delete.sort(reverse=True)

                # Delete slides
                for slide_index in slides_to_delete:
                    xml_slides = prs.slides._sldIdLst  # internal list of slide IDs
                    slides = list(xml_slides)
                    prs.part.drop_rel(slides[slide_index].rId)  # Drop the relationship
                    del xml_slides[slide_index] 
                prs.save('AdScore Reporting Template.pptx')
        except Exception as e:
            print(f"An error occurred at slide 20:  {e}") 




    return ppt_file



